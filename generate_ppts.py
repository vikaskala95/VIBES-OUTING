"""
Generate role-specific PowerPoint presentations for VIBES@Outing.
Each PPT is tailored to the role's focus area with branded slides.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Brand Colors (matching project CSS: --primary #6C3CE1) ──
BRAND_PRIMARY = RGBColor(0x6C, 0x3C, 0xE1)   # Purple primary from project
BRAND_PRIMARY_LIGHT = RGBColor(0x8E, 0x64, 0xFF)
BRAND_DARK = RGBColor(0x1E, 0x1B, 0x2E)       # #1E1B2E from project
BRAND_ORANGE = RGBColor(0xFF, 0x6B, 0x35)
BRAND_PURPLE = RGBColor(0x6C, 0x3C, 0xE1)
BRAND_TEAL = RGBColor(0x00, 0xCC, 0xA3)
BRAND_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BRAND_LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)
BRAND_GRAY = RGBColor(0x6C, 0x75, 0x7D)
BRAND_GOLD = RGBColor(0xFF, 0xD9, 0x3D)
BRAND_RED = RGBColor(0xE7, 0x4C, 0x3C)
BRAND_GREEN = RGBColor(0x2E, 0xCC, 0x71)

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "presentations_v2")
os.makedirs(OUTPUT_DIR, exist_ok=True)

# Logo image paths
LOGO_SQUARE = os.path.join(os.path.dirname(__file__), "presentations", "logo_square.png")
LOGO_WIDE = os.path.join(os.path.dirname(__file__), "presentations", "logo_wide.png")
LOGO_SMALL = os.path.join(os.path.dirname(__file__), "presentations", "logo_small.png")

WEBSITE_URL = "www.vibesouting.in"
EMAIL = "vibesoutingsupport@gmail.com"


def add_brand_header(slide, prs):
    """Add Vibes@Outing logo/brand bar to top of every slide."""
    # Top brand bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.7))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BRAND_DARK
    bar.line.fill.background()

    # Logo image from project
    if os.path.exists(LOGO_SMALL):
        slide.shapes.add_picture(LOGO_SMALL, Inches(0.2), Inches(0.08), Inches(0.55), Inches(0.55))

    # Brand name
    txBox = slide.shapes.add_textbox(Inches(0.85), Inches(0.12), Inches(3), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "VIBES@Outing"
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = BRAND_PRIMARY
    run2 = p.add_run()
    run2.text = "  •  Premium GenZ Group Outings"
    run2.font.size = Pt(10)
    run2.font.color.rgb = BRAND_GRAY


def add_brand_footer(slide, prs):
    """Add footer with tagline."""
    footer_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.8), prs.slide_width, Inches(0.7))
    footer_bar.fill.solid()
    footer_bar.fill.fore_color.rgb = BRAND_DARK
    footer_bar.line.fill.background()

    # Small logo in footer
    if os.path.exists(LOGO_SMALL):
        slide.shapes.add_picture(LOGO_SMALL, Inches(0.2), Inches(6.88), Inches(0.38), Inches(0.38))

    txBox = slide.shapes.add_textbox(Inches(0.65), Inches(6.85), Inches(6), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "VIBES@Outing — Discover. Connect. Explore Together."
    run.font.size = Pt(9)
    run.font.color.rgb = BRAND_PRIMARY

    txBox2 = slide.shapes.add_textbox(Inches(7), Inches(6.85), Inches(3), Inches(0.4))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    run2 = p2.add_run()
    run2.text = WEBSITE_URL
    run2.font.size = Pt(9)
    run2.font.color.rgb = BRAND_TEAL


def make_title_slide(prs, title, subtitle, role_tag):
    """Create a branded title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Full background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BRAND_DARK
    bg.line.fill.background()

    # Accent stripe
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.8), prs.slide_width, Inches(0.08))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = BRAND_PRIMARY
    stripe.line.fill.background()

    # Logo image (from project)
    if os.path.exists(LOGO_SQUARE):
        slide.shapes.add_picture(LOGO_SQUARE, Inches(3.85), Inches(0.5), Inches(2.3), Inches(2.3))

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(8.4), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "VIBES@Outing"
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = BRAND_PRIMARY

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.9), Inches(8), Inches(0.6))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = subtitle
    run2.font.size = Pt(20)
    run2.font.color.rgb = BRAND_WHITE

    # Role tag
    txBox3 = slide.shapes.add_textbox(Inches(1), Inches(4.7), Inches(8), Inches(0.5))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = role_tag
    run3.font.size = Pt(14)
    run3.font.color.rgb = BRAND_TEAL
    run3.font.italic = True

    # Tagline at bottom
    txBox4 = slide.shapes.add_textbox(Inches(1), Inches(5.8), Inches(8), Inches(0.5))
    tf4 = txBox4.text_frame
    p4 = tf4.paragraphs[0]
    p4.alignment = PP_ALIGN.CENTER
    run4 = p4.add_run()
    run4.text = "Discover. Connect. Explore Together."
    run4.font.size = Pt(16)
    run4.font.color.rgb = BRAND_GOLD
    run4.font.italic = True

    # Year
    txBox5 = slide.shapes.add_textbox(Inches(3), Inches(6.4), Inches(4), Inches(0.4))
    tf5 = txBox5.text_frame
    p5 = tf5.paragraphs[0]
    p5.alignment = PP_ALIGN.CENTER
    run5 = p5.add_run()
    run5.text = "Founded 2026 | India-first | GenZ-focused"
    run5.font.size = Pt(10)
    run5.font.color.rgb = BRAND_GRAY


def make_content_slide(prs, title, bullets, accent_color=BRAND_PRIMARY):
    """Standard content slide with branded header/footer."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Light bg
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BRAND_LIGHT_BG
    bg.line.fill.background()

    add_brand_header(slide, prs)
    add_brand_footer(slide, prs)

    # Accent left bar
    left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.7), Inches(0.06), Inches(6.1))
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = accent_color
    left_bar.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(0.9), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = BRAND_DARK

    # Underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.5), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = accent_color
    line.line.fill.background()

    # Bullets
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(9), Inches(4.9))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.space_after = Pt(6)
        p.space_before = Pt(2)

        if bullet.startswith("##"):
            # Sub-header
            run = p.add_run()
            run.text = bullet.replace("##", "").strip()
            run.font.size = Pt(18)
            run.font.bold = True
            run.font.color.rgb = accent_color
        elif bullet.startswith(">>"):
            # Highlighted callout
            run = p.add_run()
            run.text = "  ★  " + bullet.replace(">>", "").strip()
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = BRAND_PURPLE
        else:
            run = p.add_run()
            run.text = "●  " + bullet
            run.font.size = Pt(13)
            run.font.color.rgb = BRAND_DARK

    return slide


def make_table_slide(prs, title, headers, rows, accent_color=BRAND_PRIMARY):
    """Slide with a styled table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BRAND_LIGHT_BG
    bg.line.fill.background()

    add_brand_header(slide, prs)
    add_brand_footer(slide, prs)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(0.9), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = BRAND_DARK

    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_width = Inches(9)
    left = Inches(0.5)
    top = Inches(1.7)
    height = Inches(4.8)

    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, table_width, height)
    table = table_shape.table

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = accent_color
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(11)
                run.font.bold = True
                run.font.color.rgb = BRAND_WHITE

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BRAND_WHITE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF5)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)
                    run.font.color.rgb = BRAND_DARK

    return slide


def make_kpi_slide(prs, title, kpis, accent_color=BRAND_PRIMARY):
    """Slide with big KPI cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BRAND_LIGHT_BG
    bg.line.fill.background()

    add_brand_header(slide, prs)
    add_brand_footer(slide, prs)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(0.9), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = BRAND_DARK

    colors = [BRAND_ORANGE, BRAND_PURPLE, BRAND_TEAL, BRAND_GOLD, BRAND_GREEN, BRAND_RED]
    cols = min(len(kpis), 3)
    rows_count = (len(kpis) + cols - 1) // cols

    card_w = Inches(2.8)
    card_h = Inches(1.8)
    start_x = Inches(0.5)
    start_y = Inches(1.7)
    gap_x = Inches(0.2)
    gap_y = Inches(0.2)

    for idx, (label, value) in enumerate(kpis):
        r = idx // cols
        c = idx % cols
        x = start_x + c * (card_w + gap_x)
        y = start_y + r * (card_h + gap_y)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = BRAND_WHITE
        card.line.color.rgb = colors[idx % len(colors)]
        card.line.width = Pt(2)

        # Top color stripe
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.06))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = colors[idx % len(colors)]
        stripe.line.fill.background()

        # Value
        vBox = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.2), card_w - Inches(0.3), Inches(0.9))
        vtf = vBox.text_frame
        vtf.word_wrap = True
        vp = vtf.paragraphs[0]
        vp.alignment = PP_ALIGN.CENTER
        vr = vp.add_run()
        vr.text = str(value)
        vr.font.size = Pt(24)
        vr.font.bold = True
        vr.font.color.rgb = colors[idx % len(colors)]

        # Label
        lBox = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(1.05), card_w - Inches(0.3), Inches(0.6))
        ltf = lBox.text_frame
        ltf.word_wrap = True
        lp = ltf.paragraphs[0]
        lp.alignment = PP_ALIGN.CENTER
        lr = lp.add_run()
        lr.text = label
        lr.font.size = Pt(11)
        lr.font.color.rgb = BRAND_GRAY

    return slide


def make_closing_slide(prs, role_name):
    """Branded closing/thank you slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BRAND_DARK
    bg.line.fill.background()

    # Logo image (from project)
    if os.path.exists(LOGO_SQUARE):
        slide.shapes.add_picture(LOGO_SQUARE, Inches(3.85), Inches(0.8), Inches(2.3), Inches(2.3))

    # Thank you
    txBox = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Thank You"
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = BRAND_PRIMARY

    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.0), Inches(8), Inches(0.5))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = f"VIBES@Outing — {role_name} Deck"
    run2.font.size = Pt(16)
    run2.font.color.rgb = BRAND_WHITE

    # Contact
    txBox3 = slide.shapes.add_textbox(Inches(1.5), Inches(5.0), Inches(7), Inches(1.2))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    contacts = [
        f"🌐  {WEBSITE_URL}",
        f"📧  {EMAIL}",
        "💻  github.com/vikaskala95/vibes-outing",
    ]
    for i, c in enumerate(contacts):
        if i == 0:
            p = tf3.paragraphs[0]
        else:
            p = tf3.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = c
        run.font.size = Pt(12)
        run.font.color.rgb = BRAND_TEAL

    # Tagline
    txBox4 = slide.shapes.add_textbox(Inches(1), Inches(6.3), Inches(8), Inches(0.4))
    tf4 = txBox4.text_frame
    p4 = tf4.paragraphs[0]
    p4.alignment = PP_ALIGN.CENTER
    run4 = p4.add_run()
    run4.text = '"Touch grass. Make memories. No cap." ✨'
    run4.font.size = Pt(14)
    run4.font.italic = True
    run4.font.color.rgb = BRAND_GOLD


# ═══════════════════════════════════════════════════════════════
# 1. FOUNDER & CEO
# ═══════════════════════════════════════════════════════════════
def create_founder_ceo_ppt():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    make_title_slide(prs, "VIBES@Outing", "Founder & CEO Pitch Deck", "Strategic Vision | Investor Relations | Growth")

    make_content_slide(prs, "The Vision", [
        "##Our Mission",
        "Make every weekend an adventure for GenZ — no planning, no hassle, all vibes",
        "Bridge the loneliness gap: 64% of GenZ want to explore but lack a squad",
        "Build India's #1 GenZ social outing platform",
        "##Why Now?",
        "₹12,000 Cr adventure tourism market in India growing at 15% CAGR",
        "GenZ (18-28) is the largest demographic — digital native, experience-hungry",
        "No single platform solves discovery + booking + social for group outings",
        ">>First-mover advantage in GenZ-focused premium group outing space",
    ])

    make_content_slide(prs, "The Problem We Solve", [
        "##Pain Points",
        "Loneliness Epidemic — 64% of GenZ want to explore but have no squad",
        "No Platform — No single app for discover → join → pay for group outings",
        "Trust Gap — Going with strangers feels unsafe (no verification, no reviews)",
        "Coordination Chaos — Planning via WhatsApp/Instagram is messy",
        "Overpriced Solo Travel — Group sharing reduces cost by 40-60%",
        "##Our Solution",
        "Premium all-inclusive trips: High-end Resort + Private Cab in every outing",
        "Token-based booking: Pay just 20% upfront (₹300-₹2,600)",
        "Solo-friendly matching + Group booking in one click",
        ">>Scroll & Pick → Pay & Lock → Show Up → Slay!",
    ])

    make_kpi_slide(prs, "Key Metrics & Traction", [
        ("Curated Outings", "42"),
        ("Price Range", "₹1,499 - ₹12,999"),
        ("Avg Booking Value", "₹7,500"),
        ("Year 1 Revenue Target", "₹93L"),
        ("Year 2 Target", "₹4 Cr"),
        ("Year 3 Target", "₹12 Cr"),
    ])

    make_table_slide(prs, "Market Opportunity — TAM/SAM/SOM", 
        ["Metric", "Value", "Description"],
        [
            ["TAM", "₹12,000 Cr", "India adventure & experiential tourism"],
            ["SAM", "₹3,000 Cr", "GenZ weekend group outings in South India"],
            ["SOM", "₹75 Cr", "Year 3 target with 5 cities"],
            ["Early Adopters", "3L+ GenZ", "Bangalore alone"],
        ])

    make_table_slide(prs, "Revenue Streams", 
        ["Stream", "Model", "Year 1 Projected"],
        [
            ["Booking Commission", "15-20% per booking", "₹45L"],
            ["Premium Membership", "₹999/year", "₹12L"],
            ["Featured Listings", "Resort/organizer promotions", "₹10L"],
            ["Corporate Packages", "Startup team outings", "₹18L"],
            ["Affiliate Earnings", "Travel partners", "₹8L"],
        ])

    make_content_slide(prs, "Competitive Moat", [
        "##What Makes Us Un-copyable",
        "GenZ-only community with verified users (Aadhaar/PAN)",
        "35+ curated premium all-inclusive trips (Resort + Cab included)",
        "Token-based booking (20% upfront) — unique in market",
        "AI-powered interest matching + MCP server integration",
        "Solo → Group matching engine — no competitor does this well",
        "##Competitive Advantages over:",
        "MakeMyTrip: No group matching, not GenZ-focused",
        "Meetup.com: No booking/payments, no premium trips",
        "WhatsApp Groups: No trust layer, no logistics, messy coordination",
        ">>PWA + MCP AI integration = tech moat",
    ])

    make_content_slide(prs, "Growth Roadmap", [
        "##2026 — Foundation (Current)",
        "MVP live: 42 curated outings, Razorpay payments, PWA, MCP server",
        "Bangalore launch with college & startup communities",
        "##2027 — Scale",
        "5 cities: Hyderabad, Chennai, Pune, Mumbai",
        "1.5L users, 100+ outings, Mobile app (React Native)",
        "##2028 — Dominate",
        "Pan-India, AI-powered social matching",
        "Organizer self-service portal, trip albums & social features",
        "##2029-2030 — Expand",
        "Southeast Asia expansion",
        ">>Goal: #1 GenZ social outing platform in Asia by 2030",
    ])

    make_content_slide(prs, "Funding Ask — Seed Round: ₹30 Lakhs", [
        "##Use of Funds",
        "Product Development (Mobile App): 35% — ₹10.5L",
        "Marketing & User Acquisition: 35% — ₹10.5L",
        "Operations & Team: 20% — ₹6L",
        "Legal & Compliance: 10% — ₹3L",
        "##Deal Terms",
        "12-month runway",
        "Target: 25,000 users, 5 cities, ₹55L revenue",
        "Equity Offered: 8-12%",
        ">>Looking for strategic investors with travel/GenZ ecosystem access",
    ])

    make_content_slide(prs, "Team & Leadership", [
        "##Core Team Roles",
        "Founder/CEO — Vision, strategy, fundraising, partnerships",
        "Full-Stack Developer — Platform (Node.js/Express, PostgreSQL, Razorpay)",
        "Product Manager — Feature roadmap, user research, outing curation",
        "UI/UX Designer — Mobile-first GenZ design, PWA experience",
        "Marketing Lead — User acquisition, influencer outreach, SEO",
        "Operations Lead — Vendor management, resort partnerships, logistics",
        "Community Manager — Engagement, reviews moderation, NPS tracking",
        "Finance & Compliance — Payments, taxes, legal, privacy",
    ])

    make_closing_slide(prs, "Founder & CEO")
    prs.save(os.path.join(OUTPUT_DIR, "01_Founder_CEO_Deck.pptx"))
    print("✅ Founder & CEO deck created")


# ═══════════════════════════════════════════════════════════════
# 2. MARKETING MANAGER
# ═══════════════════════════════════════════════════════════════
def create_marketing_manager_ppt():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    make_title_slide(prs, "VIBES@Outing", "Marketing Manager Playbook", "User Acquisition | Brand Strategy | Growth Hacking")

    make_content_slide(prs, "Marketing Vision", [
        "##Brand Positioning",
        "VIBES@Outing = The GenZ Squad Outing App",
        "Tone: Fun, authentic, FOMO-inducing, meme-friendly",
        "Target: 18-28 year olds — adventure seekers, Insta-lovers, solo travelers",
        "##Brand Promise",
        "Zero planning. Premium trips. All vibes.",
        "Every stranger can become a friend",
        ">>42 curated trips from ₹1,499 to ₹12,999 — all-inclusive",
    ], BRAND_PURPLE)

    make_content_slide(prs, "Target Audience Personas", [
        "##Persona 1: Solo Explorer (35% of users)",
        "Age 22-26, working professional, wants to travel but friends aren't free",
        "Values safety (ID verification), affordability (group pricing)",
        "##Persona 2: College Squad Leader (30%)",
        "Age 18-22, organizes group trips, hates coordination chaos",
        "Values one-click group booking, premium experiences",
        "##Persona 3: Weekend Warrior (20%)",
        "Age 24-28, young professional, seeks regular adventures",
        "Values variety (42 destinations), premium membership perks",
        "##Persona 4: Insta-Aesthetic Traveler (15%)",
        "Age 20-25, content creator, seeks photo-worthy experiences",
        "Values curated destinations, trip galleries, social sharing",
    ], BRAND_PURPLE)

    make_content_slide(prs, "Go-to-Market Strategy", [
        "##Phase 1 — Bangalore Launch (Month 1-3)",
        "Instagram Reels + College WhatsApp groups + Reddit r/bangalore",
        "GenZ influencer partnerships (travel micro-influencers, 10K-100K followers)",
        "₹0 viral referral system: Invite friend = ₹500 off next trip",
        "College campus ambassadors — 10 major colleges in Bangalore",
        "##Phase 2 — South India (Month 4-8)",
        "Expand to Hyderabad, Chennai, Pune — localized content",
        "Startup team outing partnerships (corporate GenZ packages)",
        "##Phase 3 — Scale (Month 9-12)",
        "Paid ads: Instagram, YouTube Shorts, Google Search",
        "Brand collaborations with travel gear & lifestyle brands",
        ">>Goal: 25,000 users by end of Year 1",
    ], BRAND_PURPLE)

    make_table_slide(prs, "Channel Strategy & Budget", 
        ["Channel", "Strategy", "Budget %", "Expected CAC"],
        [
            ["Instagram Reels/Stories", "3x/week trip highlights, UGC, behind-scenes", "25%", "₹80"],
            ["College WhatsApp/Telegram", "Campus ambassadors, group promotions", "15%", "₹30"],
            ["Referral Program", "₹500 off for inviter + invitee", "20%", "₹50"],
            ["Google SEO/SEM", "Long-tail: 'weekend trips from Bangalore'", "15%", "₹120"],
            ["YouTube Shorts", "Trip vlogs, destination reveals", "10%", "₹100"],
            ["Influencer Collabs", "Micro-influencers (10K-100K)", "10%", "₹90"],
            ["Reddit/Twitter", "Community engagement, trip AMAs", "5%", "₹40"],
        ], BRAND_PURPLE)

    make_kpi_slide(prs, "Marketing KPIs & Targets", [
        ("Conversion Rate", "8%"),
        ("CAC Target", "< ₹100"),
        ("Repeat Booking", "35%"),
        ("Referral Rate", "25%"),
        ("NPS Score", "60+"),
        ("Monthly Active Users (Q4)", "25K"),
    ])

    make_content_slide(prs, "Content Marketing Strategy", [
        "##Content Pillars",
        "Trip Highlights — Reels/Shorts from each of 42 destinations",
        "User Stories — Solo travelers who found their squad",
        "Behind-the-Scenes — Resort partnerships, team culture",
        "FOMO Content — 'Last 3 spots left!' urgency posts",
        "##Content Calendar",
        "Monday: Destination Monday (new trip spotlight)",
        "Wednesday: Squad Stories (user testimonials)",
        "Friday: Weekend FOMO (upcoming trip promos)",
        "Sunday: Trip Recap (photo dumps from completed outings)",
        "##SEO Strategy",
        "Blog posts: 'Best weekend trips from Bangalore', 'Solo travel guide GenZ'",
        ">>Target 50,000 monthly organic visitors by Month 6",
    ], BRAND_PURPLE)

    make_content_slide(prs, "Referral & Loyalty Program", [
        "##Viral Referral Engine",
        "Refer a friend → Both get ₹500 off next trip",
        "3 referrals → Free day trip (₹1,499 value)",
        "10 referrals → Premium membership free for 1 year",
        "##Premium Membership (₹999/year)",
        "Early access to new trips (24hr head-start)",
        "Exclusive members-only outings",
        "5% discount on all bookings",
        "Priority group matching",
        "##Gamification",
        "Badges: Explorer, Adventurer, Trailblazer, Legend",
        "Leaderboard for most trips booked",
        ">>Target: 25% of users activate referral within first month",
    ], BRAND_PURPLE)

    make_table_slide(prs, "Growth Projections", 
        ["Quarter", "Users", "Monthly Outings", "Revenue", "Marketing Spend"],
        [
            ["Q1 (Launch)", "500", "15", "₹8L", "₹2L"],
            ["Q2", "3,000", "35", "₹20L", "₹4L"],
            ["Q3", "10,000", "60", "₹35L", "₹5L"],
            ["Q4", "25,000", "100", "₹55L", "₹6L"],
        ], BRAND_PURPLE)

    make_content_slide(prs, "Brand Identity & Guidelines", [
        "##Visual Identity",
        "Primary Color: Orange (#FF6B35) — Energy, adventure, excitement",
        "Secondary: Dark Navy (#1A1A2E) — Trust, premium feel",
        "Accent: Teal (#00CCA3) — Fresh, GenZ-friendly",
        "##Tone of Voice",
        "Casual, fun, meme-aware — never corporate",
        "Use GenZ slang naturally: 'no cap', 'slay', 'vibes', 'touch grass'",
        "Emojis encouraged in social media",
        "##Brand Taglines",
        "Primary: 'Discover. Connect. Explore Together.'",
        "Social: 'Touch grass. Make memories. No cap.'",
        ">>Consistency across all 42 destination pages + social channels",
    ], BRAND_PURPLE)

    make_closing_slide(prs, "Marketing Manager")
    prs.save(os.path.join(OUTPUT_DIR, "02_Marketing_Manager_Deck.pptx"))
    print("✅ Marketing Manager deck created")


# ═══════════════════════════════════════════════════════════════
# 3. DIGITAL MEDIA MARKETING MANAGER
# ═══════════════════════════════════════════════════════════════
def create_digital_media_ppt():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    make_title_slide(prs, "VIBES@Outing", "Digital Media Marketing Deck", "Social Media | Content | Influencer | Performance Ads")

    make_content_slide(prs, "Digital Media Landscape", [
        "##Where GenZ Lives Online",
        "Instagram: 85% of GenZ India — primary discovery platform",
        "YouTube Shorts: 70% daily consumption — trip vlogs & destination reveals",
        "Reddit: r/bangalore (500K+) — organic community engagement",
        "WhatsApp/Telegram: Group discovery & viral sharing",
        "##Our Digital Advantage",
        "42 curated photogenic destinations = unlimited content goldmine",
        "User-generated content from every trip = authentic social proof",
        ">>Each trip = 15-20 Instagram-worthy moments for UGC",
    ], BRAND_TEAL)

    make_table_slide(prs, "Social Media Platform Strategy", 
        ["Platform", "Content Type", "Frequency", "Goal"],
        [
            ["Instagram Reels", "30-sec trip highlights, transitions", "5x/week", "Discovery + FOMO"],
            ["Instagram Stories", "Polls, countdowns, behind-scenes", "Daily", "Engagement"],
            ["Instagram Posts", "Carousels, trip galleries, reviews", "3x/week", "Trust + SEO"],
            ["YouTube Shorts", "Trip vlogs, 'Day in the life' of outing", "3x/week", "Reach expansion"],
            ["YouTube Long-form", "Full trip documentaries, reviews", "2x/month", "SEO + depth"],
            ["Reddit", "Trip reports, AMAs, community Q&A", "3x/week", "Organic trust"],
            ["Twitter/X", "Memes, trending takes, quick updates", "Daily", "Virality"],
            ["WhatsApp Channels", "Trip alerts, exclusive offers", "2x/week", "Direct conversion"],
        ], BRAND_TEAL)

    make_content_slide(prs, "Content Production Pipeline", [
        "##Content Types (Priority Order)",
        "Reels/Shorts: Trip highlights with trending audio (highest ROI)",
        "User testimonials: Solo travelers who found their squad",
        "Destination carousels: '5 reasons to visit Coorg this weekend'",
        "Behind-the-scenes: Resort setup, cab arrangements, team prep",
        "FOMO triggers: 'Only 3 spots left!' countdown stories",
        "##Production Workflow",
        "Pre-trip: Teaser content, countdown, destination spotlight",
        "During trip: Live stories, quick snaps, attendee takeovers",
        "Post-trip: Photo dumps, review videos, user UGC reposts",
        ">>Target: 100+ unique content pieces per month across platforms",
    ], BRAND_TEAL)

    make_content_slide(prs, "Influencer Marketing Strategy", [
        "##Tier 1: Micro-Influencers (10K-50K followers)",
        "30 travel/lifestyle creators in Bangalore — barter + small fee",
        "Deliverables: 1 Reel + 3 Stories per trip attended",
        "Expected reach: 300K-500K per month",
        "##Tier 2: Mid-Tier (50K-200K followers)",
        "5-8 curated partnerships per quarter",
        "Paid collaboration: ₹15K-₹30K per campaign",
        "##Tier 3: Macro/Celebrity (200K+)",
        "1-2 brand ambassadors for annual deals (Phase 2)",
        "##Creator Program",
        "Free trips for creators who produce 3+ pieces of content",
        "Exclusive 'Creator Squad' badge on platform",
        ">>Influencer content typically generates 3x more trust than brand content",
    ], BRAND_TEAL)

    make_table_slide(prs, "Paid Advertising Strategy", 
        ["Platform", "Ad Type", "Target", "Monthly Budget", "Expected ROAS"],
        [
            ["Instagram Ads", "Reel ads, carousel ads", "18-28, Bangalore, interests: travel", "₹50K", "3x"],
            ["Google Search", "Long-tail keywords", "'weekend trips Bangalore', 'group outing'", "₹30K", "4x"],
            ["YouTube Pre-roll", "6-sec bumpers", "18-28, travel intent", "₹20K", "2.5x"],
            ["Facebook/Meta", "Lookalike audiences", "Based on existing user data", "₹15K", "2x"],
            ["Reddit Promoted", "Promoted posts in r/bangalore", "Travel + GenZ subreddits", "₹10K", "3.5x"],
        ], BRAND_TEAL)

    make_kpi_slide(prs, "Digital Media KPIs", [
        ("Instagram Followers (Y1)", "50K"),
        ("Monthly Reach", "500K"),
        ("Engagement Rate", "> 5%"),
        ("Content-to-Booking Rate", "3%"),
        ("Cost per Click", "< ₹5"),
        ("ROAS Target", "3x+"),
    ])

    make_content_slide(prs, "SEO & Content Marketing", [
        "##Blog/SEO Strategy",
        "Target keywords: 'weekend trips from Bangalore', 'GenZ group travel'",
        "'Best of' listicles for each of 42 destinations",
        "Trip planning guides with embedded booking CTAs",
        "##SEO Technical",
        "PWA already optimized for mobile-first indexing",
        "Structured data for trip listings (schema.org/Event)",
        "Page speed optimization — current PWA loads in < 2 seconds",
        "##Email Marketing",
        "Welcome series: 3-email onboarding flow",
        "Weekly newsletter: 'This Weekend's Top Picks'",
        "Abandoned booking recovery emails",
        ">>Target: 50,000 monthly organic visitors by Month 6",
    ], BRAND_TEAL)

    make_content_slide(prs, "Viral & Guerrilla Campaigns", [
        "##Campaign Ideas",
        "'Mystery Trip' — Book without knowing destination, revealed day-of",
        "'Stranger Squad Challenge' — Film making friends on VIBES@Outing trips",
        "'Weekend Bingo' — Visit 5 destinations, win free Goa trip",
        "'Touch Grass Challenge' — TikTok/Reels challenge, most views = free trip",
        "##Partnerships",
        "College festivals: Sponsor travel-themed events at top 10 Bangalore colleges",
        "Startup team outing partnerships: Offer corporate GenZ packages",
        "Travel gear brands: Co-branded content with backpack/clothing brands",
        "##Community Building",
        "VIBES@Outing Discord/Telegram community for trip discussions",
        "Monthly meetups in Bangalore for platform users",
        ">>Every campaign ties back to app installs + first booking",
    ], BRAND_TEAL)

    make_closing_slide(prs, "Digital Media Marketing Manager")
    prs.save(os.path.join(OUTPUT_DIR, "03_Digital_Media_Marketing_Deck.pptx"))
    print("✅ Digital Media Marketing Manager deck created")


# ═══════════════════════════════════════════════════════════════
# 4. HEAD OF OPERATIONS
# ═══════════════════════════════════════════════════════════════
def create_operations_ppt():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    make_title_slide(prs, "VIBES@Outing", "Head of Operations Deck", "Logistics | Vendor Management | Trip Execution | Safety")

    make_content_slide(prs, "Operations Overview", [
        "##Mission",
        "Flawless execution of 42 curated premium trips from Bangalore",
        "Every trip includes: High-end Resort + Private Cab (non-negotiable)",
        "##Scale",
        "42 destinations across 6 categories: Mountains, Adventure, Beaches, Wildlife, Heritage, Road Trips",
        "14 Day Trips + 21 Overnight Trips + 5 Wildlife/Safari + 2 Special",
        "Price range: ₹1,499 to ₹12,999 per person",
        "Average group size: 15 people",
        "##Year 1 Targets",
        "Q1: 15 outings/month → Q4: 100 outings/month",
        ">>Zero safety incidents. 100% on-time departures.",
    ])

    make_table_slide(prs, "Trip Categories & Operations", 
        ["Category", "Count", "Avg Duration", "Key Logistics", "Complexity"],
        [
            ["Day Trips", "14", "8-12 hours", "Cab + activity + lunch", "Low"],
            ["2-Day/1-Night", "21", "2 days", "Resort + cab + meals + activities", "Medium"],
            ["Wildlife Safari", "5", "2 days", "Resort + safari permits + guides", "High"],
            ["Adventure", "6", "1-2 days", "Safety gear + trained guides + insurance", "High"],
            ["Heritage", "4", "1-2 days", "Guides + temple timings + cultural etiquette", "Medium"],
            ["Beach/Coastal", "5", "2 days", "Water sports vendors + safety equipment", "Medium"],
        ])

    make_content_slide(prs, "Vendor & Partner Management", [
        "##Resort Partners",
        "Tier-1 resorts at each of 42 destinations — quality audit quarterly",
        "Negotiate bulk rates: 30-40% below rack rate for consistent bookings",
        "Backup resort at every destination for contingency",
        "##Cab/Transport Partners",
        "3 cab partners minimum — Tempo Traveller (12-15 seats) + Mini Bus (20+)",
        "GPS tracking on all vehicles, driver verification (DL + background check)",
        "##Activity Vendors",
        "Certified adventure operators: rafting, trekking, rappelling, kayaking",
        "Safety certifications mandatory, insurance verification",
        "##Food & Beverage",
        "Pre-negotiated meal packages at destination restaurants",
        "Dietary preference collection during booking (veg/non-veg/vegan)",
        ">>All vendors rated after every trip — below 4.0 rating = review/replace",
    ])

    make_content_slide(prs, "Trip Execution Playbook", [
        "##T-7 Days (Pre-Trip)",
        "Confirm resort booking, cab assignment, activity slots",
        "Send itinerary to all booked participants via email + WhatsApp",
        "Collect emergency contact, dietary preferences, medical info",
        "##T-1 Day",
        "Final headcount confirmation, collect remaining 80% payment",
        "Cab driver briefing: route, stops, contact numbers",
        "Weather check — contingency plan activation if needed",
        "##Trip Day",
        "Trip coordinator on-ground for every outing (15+ group)",
        "Real-time check-in/check-out tracking",
        "Emergency protocol: Local hospital, police, insurance contacts ready",
        "##T+1 Day (Post-Trip)",
        "Send feedback survey to all attendees",
        "Process ratings & reviews, photo gallery upload",
        ">>Vendor settlement within 48 hours of trip completion",
    ])

    make_content_slide(prs, "Safety & Emergency Protocols", [
        "##Safety Standards",
        "ID verification mandatory: Aadhaar/PAN before trip",
        "Emergency contacts collected for every participant",
        "First-aid kit on every cab, trained coordinator for groups 15+",
        "##Emergency Playbook",
        "Medical emergency: Nearest hospital pre-mapped for all 42 destinations",
        "Vehicle breakdown: Backup cab dispatch within 2 hours (partner SLA)",
        "Weather emergency: Alternate indoor activity plan for each destination",
        "##Insurance",
        "Group travel insurance for every trip (included in price)",
        "Adventure activity insurance: Separate coverage for rafting/trekking/rappelling",
        ">>Zero tolerance for unverified participants on any trip",
    ])

    make_kpi_slide(prs, "Operations KPIs", [
        ("On-Time Departure Rate", "98%+"),
        ("Vendor Quality Score", "> 4.0/5"),
        ("Safety Incidents", "Zero"),
        ("Trip NPS", "60+"),
        ("Vendor Settlement Time", "< 48 hrs"),
        ("Outing Capacity Utilization", "> 80%"),
    ])

    make_content_slide(prs, "Scaling Operations", [
        "##Phase 1: Bangalore (Month 1-3)",
        "15→100 outings/month ramp-up, 42 destinations operational",
        "Core team: 2 operations coordinators + 5 trip leaders",
        "##Phase 2: Multi-City (Month 4-8)",
        "Hyderabad, Chennai, Pune — replicate Bangalore playbook",
        "Regional operations leads in each city",
        "Expand to 75+ outings in catalog",
        "##Phase 3: Organizer Platform (Month 9-12)",
        "Self-service portal for verified organizers to create trips",
        "Operations quality audit for third-party trips",
        "Automated vendor management system",
        ">>Standard Operating Procedures documented for all 42 destinations",
    ])

    make_closing_slide(prs, "Head of Operations")
    prs.save(os.path.join(OUTPUT_DIR, "04_Head_of_Operations_Deck.pptx"))
    print("✅ Head of Operations deck created")


# ═══════════════════════════════════════════════════════════════
# 5. FINANCE MANAGER
# ═══════════════════════════════════════════════════════════════
def create_finance_manager_ppt():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    make_title_slide(prs, "VIBES@Outing", "Finance Manager Deck", "Revenue | Unit Economics | Compliance | Financial Planning")

    make_content_slide(prs, "Financial Overview", [
        "##Business Model",
        "Marketplace model: Commission on bookings + membership + promotions",
        "Token-based booking: 20% upfront, 80% within 24hrs of trip date",
        "Payment via Razorpay: UPI, Cards, NetBanking",
        "##Revenue Streams (5 pillars)",
        "Booking Commission: 15-20% per booking",
        "Premium Membership: ₹999/year",
        "Featured Listings: Resort/organizer promotions",
        "Corporate Packages: Startup team outings",
        "Affiliate Earnings: Travel partners",
        ">>Year 1 Total Revenue Target: ₹93 Lakhs",
    ], BRAND_GREEN)

    make_table_slide(prs, "Revenue Projections", 
        ["Period", "Users", "Monthly Outings", "Avg Booking", "Revenue"],
        [
            ["Q1", "500", "15", "₹7,500", "₹8L"],
            ["Q2", "3,000", "35", "₹7,500", "₹20L"],
            ["Q3", "10,000", "60", "₹7,500", "₹35L"],
            ["Q4", "25,000", "100", "₹7,500", "₹55L"],
            ["Year 1 Total", "25,000", "210 total", "₹7,500", "₹93L"],
            ["Year 2", "1,50,000", "3,000 total", "₹8,000", "₹4Cr"],
            ["Year 3", "5,00,000", "10,000 total", "₹8,500", "₹12Cr"],
        ], BRAND_GREEN)

    make_kpi_slide(prs, "Unit Economics", [
        ("Avg Booking Value", "₹7,500"),
        ("Commission Rate", "15-20%"),
        ("Gross Margin/Booking", "₹1,125-₹1,500"),
        ("CAC Target", "< ₹100"),
        ("LTV (3 trips)", "₹3,375-₹4,500"),
        ("LTV:CAC Ratio", "30-45x"),
    ])

    make_content_slide(prs, "Token-Based Booking Economics", [
        "##How Token Booking Works",
        "Customer pays 20% upfront to lock their spot (₹300-₹2,600)",
        "Remaining 80% due within 24 hours of trip date",
        "##Financial Benefits",
        "Cash flow positive: Token amount received 7-30 days before trip",
        "Lower refund liability: Only token amount at risk for cancellations",
        "Higher commitment: Token creates psychological lock-in",
        "##Cancellation Policy",
        "7+ days before trip: 90% token refund",
        "3-7 days: 50% token refund",
        "< 3 days: No refund (token forfeited)",
        ">>Token forfeitures estimated at 5-8% — additional revenue stream",
    ], BRAND_GREEN)

    make_table_slide(prs, "Cost Structure", 
        ["Category", "Monthly (Q1)", "Monthly (Q4)", "% of Revenue"],
        [
            ["Vendor Payments (Resorts+Cabs)", "₹5.5L", "₹35L", "60-65%"],
            ["Marketing & Acquisition", "₹0.7L", "₹2L", "8-10%"],
            ["Team Salaries (10 roles)", "₹3L", "₹5L", "12-15%"],
            ["Tech Infrastructure", "₹0.3L", "₹0.5L", "3-5%"],
            ["Payment Gateway Fees (2%)", "₹0.16L", "₹1.1L", "2%"],
            ["Legal & Compliance", "₹0.25L", "₹0.25L", "1-2%"],
            ["Contingency & Insurance", "₹0.2L", "₹0.5L", "2-3%"],
        ], BRAND_GREEN)

    make_content_slide(prs, "Seed Funding — ₹30 Lakhs", [
        "##Fund Allocation",
        "Product Development (Mobile App): 35% — ₹10.5L",
        "Marketing & User Acquisition: 35% — ₹10.5L",
        "Operations & Team: 20% — ₹6L",
        "Legal & Compliance: 10% — ₹3L",
        "##Financial Milestones",
        "Month 3: Break-even on unit economics (per trip profitable)",
        "Month 6: Operating cash flow positive",
        "Month 12: ₹55L cumulative revenue, ready for Series A",
        "##Equity Terms",
        "8-12% equity for ₹30L seed round",
        ">>12-month runway with conservative growth assumptions",
    ], BRAND_GREEN)

    make_content_slide(prs, "Payment & Compliance", [
        "##Payment Infrastructure",
        "Razorpay integration: UPI, Cards, NetBanking",
        "Automated split payments: Commission retained, vendor paid",
        "Refund automation: Policy-based refund processing",
        "##Tax Compliance",
        "GST registration & quarterly filing",
        "TDS deduction on vendor payments",
        "Income tax compliance (startup exemption under Section 80-IAC if eligible)",
        "##Legal & Data Privacy",
        "Terms of Service & Privacy Policy (DPDP Act 2023 compliant)",
        "Aadhaar/PAN data handling — encrypted storage, limited access",
        "Travel insurance compliance for group activities",
        ">>Monthly financial reporting dashboard in admin panel",
    ], BRAND_GREEN)

    make_content_slide(prs, "Financial Risk Management", [
        "##Key Risks & Mitigation",
        "Cancellation spikes: Token system limits exposure to 20% of booking value",
        "Vendor default: 2+ backup vendors per destination, advance deposits",
        "Seasonality: Off-season pricing, indoor/urban activities in monsoon",
        "##Cash Flow Management",
        "Token collected upfront → positive working capital cycle",
        "Vendor payments on net-7 terms (post-trip settlement)",
        "Membership revenue provides recurring baseline",
        "##Sensitivity Analysis",
        "If bookings 30% below target: Still cash-flow positive at Month 8",
        "If CAC doubles: Shift to organic/referral — break-even at Month 10",
        ">>Conservative model — all projections assume 70% of optimistic scenario",
    ], BRAND_GREEN)

    make_closing_slide(prs, "Finance Manager")
    prs.save(os.path.join(OUTPUT_DIR, "05_Finance_Manager_Deck.pptx"))
    print("✅ Finance Manager deck created")


# ═══════════════════════════════════════════════════════════════
# 6. CUSTOMER RELATIONSHIP MANAGER
# ═══════════════════════════════════════════════════════════════
def create_customer_relationship_ppt():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    make_title_slide(prs, "VIBES@Outing", "Customer Relationship Manager Deck", "Community | Engagement | Retention | Support")

    make_content_slide(prs, "CRM Vision", [
        "##Customer-First Philosophy",
        "Every user is a potential brand ambassador — treat them like VIPs",
        "GenZ expects instant responses, authentic communication, personal touch",
        "##CRM Goals",
        "Repeat booking rate: 35%+ (users booking 2+ trips)",
        "NPS Score: 60+ (promoter-dominant)",
        "Support response time: < 2 hours",
        "Community engagement: Active user base on Discord/Telegram",
        ">>Turn solo travelers into squad leaders, and squad leaders into brand evangelists",
    ], BRAND_GOLD)

    make_content_slide(prs, "Customer Journey Mapping", [
        "##Stage 1: Discovery",
        "User finds VIBES@Outing via Instagram/referral/search",
        "CRM action: Track source, personalize first interaction",
        "##Stage 2: First Booking",
        "User browses 42 destinations, selects trip, pays token",
        "CRM action: Welcome email, trip prep guide, WhatsApp group invite",
        "##Stage 3: Trip Experience",
        "User attends outing, interacts with group, creates memories",
        "CRM action: In-trip check-ins, photo sharing, coordinator support",
        "##Stage 4: Post-Trip",
        "User receives feedback survey, rates experience, writes review",
        "CRM action: Thank-you email, review request, next trip recommendation",
        "##Stage 5: Retention & Advocacy",
        "User books again, refers friends, shares on social media",
        ">>CRM action: Loyalty rewards, referral incentives, VIP access",
    ], BRAND_GOLD)

    make_content_slide(prs, "Community Management", [
        "##Community Channels",
        "In-App Group Chat: Auto-created for each trip (pre-trip + during)",
        "WhatsApp Groups: Trip-specific groups for coordination",
        "Discord/Telegram: Platform-wide community for trip discussions",
        "##Moderation",
        "AI-assisted content moderation for group chats",
        "Clear community guidelines: No harassment, spam, or inappropriate content",
        "Report system with 24-hour resolution SLA",
        "##Engagement Programs",
        "Monthly virtual meetup: 'VIBES@Outing Community Night'",
        "Trip photo contests: Best photo wins free trip",
        "User spotlights: Feature top travelers on Instagram",
        ">>Build emotional connection — community is our competitive moat",
    ], BRAND_GOLD)

    make_table_slide(prs, "Customer Support Framework", 
        ["Category", "Channel", "Response SLA", "Resolution SLA"],
        [
            ["Pre-booking Queries", "Chat/Email/WhatsApp", "30 mins", "2 hours"],
            ["Payment Issues", "Email + Razorpay", "1 hour", "4 hours"],
            ["Trip Changes/Cancellation", "Email/Phone", "2 hours", "24 hours"],
            ["During-Trip Emergency", "Phone (Hotline)", "Immediate", "Real-time"],
            ["Post-Trip Complaints", "Email/In-App", "4 hours", "48 hours"],
            ["Review Disputes", "Email", "24 hours", "72 hours"],
            ["Refund Requests", "Email + Dashboard", "4 hours", "5 business days"],
        ], BRAND_GOLD)

    make_kpi_slide(prs, "CRM KPIs & Targets", [
        ("Repeat Booking Rate", "35%+"),
        ("NPS Score", "60+"),
        ("Support Response Time", "< 2 hrs"),
        ("Review Completion Rate", "40%"),
        ("Referral Activation", "25%"),
        ("Churn Rate", "< 15%"),
    ])

    make_content_slide(prs, "Retention & Loyalty Strategy", [
        "##Retention Levers",
        "Post-trip recommendation engine: 'Based on your Coorg trip, try Wayanad next'",
        "Birthday/anniversary trip discounts (personal touch)",
        "Milestone rewards: 3rd trip free upgrade, 5th trip free day-trip",
        "##Loyalty Tiers",
        "Explorer: 1 trip — Welcome badge, basic perks",
        "Adventurer: 3 trips — 5% discount, early access",
        "Trailblazer: 5 trips — 10% discount, exclusive trips, priority matching",
        "Legend: 10+ trips — Free annual membership, VIP experiences, ambassador status",
        "##Win-Back Campaigns",
        "30-day inactive: 'Miss the vibes? Here's ₹300 off your next trip'",
        "60-day inactive: 'Your squad misses you — exclusive trip invitation'",
        ">>Personalization is key — use trip history, interests, group preferences",
    ], BRAND_GOLD)

    make_content_slide(prs, "Feedback & Review System", [
        "##Review Collection",
        "Automated post-trip survey (NPS + detailed feedback) — sent T+1 day",
        "In-app star ratings for: Trip, Resort, Transport, Activities, Coordinator",
        "Photo review incentive: Upload trip photos = ₹100 credit",
        "##Review Analysis",
        "Sentiment analysis on all reviews — flag negative patterns",
        "Monthly CRM report: Top complaints, praise areas, actionable improvements",
        "Vendor quality tracking: Reviews feed into vendor scorecard",
        "##User Insights",
        "Trip preference analytics: Popular destinations, price sensitivity, group size",
        "Seasonal demand patterns for operations planning",
        ">>Every review is responded to within 24 hours — positive or negative",
    ], BRAND_GOLD)

    make_closing_slide(prs, "Customer Relationship Manager")
    prs.save(os.path.join(OUTPUT_DIR, "06_Customer_Relationship_Deck.pptx"))
    print("✅ Customer Relationship Manager deck created")


# ═══════════════════════════════════════════════════════════════
# 7. TECH LEAD
# ═══════════════════════════════════════════════════════════════
def create_tech_lead_ppt():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    make_title_slide(prs, "VIBES@Outing", "Tech Lead Deck", "Architecture | Stack | Security | DevOps | Roadmap")

    make_content_slide(prs, "Technical Architecture Overview", [
        "##Current Stack (MVP — Live)",
        "Frontend: Single Page App (HTML5/CSS3/Vanilla JS) + PWA",
        "Backend: Node.js + Express.js",
        "Database: PostgreSQL (via pg Pool)",
        "Payments: Razorpay (UPI, Cards, NetBanking)",
        "Auth: JWT + bcrypt + Google OAuth",
        "AI: MCP Server (@modelcontextprotocol/sdk)",
        "##Hosting",
        "Backend + DB: Render.com",
        "Frontend: Vercel",
        "Containerization: Docker (Dockerfile + render.yaml + vercel.json)",
        ">>Full MVP live with 42 curated outings, payments, admin dashboard",
    ], BRAND_PURPLE)

    make_table_slide(prs, "Tech Stack Deep Dive", 
        ["Layer", "Technology", "Purpose", "Status"],
        [
            ["Frontend", "HTML5/CSS3/JS + PWA", "Installable mobile web app", "✅ Live"],
            ["Backend API", "Node.js + Express", "REST API, business logic", "✅ Live"],
            ["Database", "PostgreSQL", "Users, bookings, outings, reviews", "✅ Live"],
            ["Payments", "Razorpay SDK", "UPI, Cards, NetBanking", "✅ Live"],
            ["Auth", "JWT + bcrypt + Google OAuth", "Secure authentication", "✅ Live"],
            ["Email", "Nodemailer (SMTP)", "Notifications, password reset", "✅ Live"],
            ["Security", "Helmet, CORS, HPP, rate limiting", "API protection", "✅ Live"],
            ["AI/MCP", "@modelcontextprotocol/sdk", "AI agent integration", "✅ Live"],
            ["DevOps", "Docker + Render + Vercel", "CI/CD, hosting", "✅ Live"],
            ["Mobile App", "React Native", "Native mobile experience", "Phase 2"],
        ], BRAND_PURPLE)

    make_content_slide(prs, "API Architecture", [
        "##Core API Endpoints",
        "Auth: /api/auth/register, /login, /google-login, /reset-password",
        "Outings: /api/outings (CRUD), /api/outings/filter, /api/outings/search",
        "Bookings: /api/bookings/create, /verify-payment, /my-bookings",
        "Reviews: /api/reviews/create, /by-outing, /by-user",
        "Admin: /api/admin/dashboard, /users, /bookings, /analytics",
        "##API Design Principles",
        "RESTful design with consistent error handling",
        "Input validation via express-validator on all endpoints",
        "Rate limiting: 100 req/15min per IP (general), 5 req/15min (auth)",
        "JWT token-based auth with refresh token rotation",
        "##Database Design",
        "Normalized PostgreSQL schema: users, outings, bookings, reviews, suggestions",
        "Indexed queries for search & filter performance",
        ">>Connection pooling via pg Pool — handles 100+ concurrent connections",
    ], BRAND_PURPLE)

    make_content_slide(prs, "Security Architecture", [
        "##Application Security",
        "Helmet.js: HTTP security headers (CSP, X-Frame-Options, etc.)",
        "CORS: Whitelist-based origin control",
        "HPP: HTTP parameter pollution prevention",
        "Rate limiting: Tiered limits (auth stricter than general API)",
        "##Authentication Security",
        "bcrypt password hashing (12 rounds)",
        "JWT with short expiry + refresh token rotation",
        "Google OAuth as social login alternative",
        "##Data Security",
        "Input sanitization & validation on all endpoints",
        "SQL injection prevention via parameterized queries (pg Pool)",
        "XSS prevention via Content Security Policy",
        "Aadhaar/PAN data: Encrypted at rest, access-controlled",
        "##Payment Security",
        "Razorpay handles PCI-DSS compliance",
        ">>Server-side payment verification — never trust client-side",
    ], BRAND_PURPLE)

    make_content_slide(prs, "PWA & Performance", [
        "##Progressive Web App",
        "Service Worker: Offline caching, background sync",
        "manifest.json: Installable on mobile home screen",
        "Push notifications (planned)",
        "##Performance Targets",
        "First Contentful Paint: < 1.5 seconds",
        "Time to Interactive: < 3 seconds",
        "Lighthouse Score: 90+ across all categories",
        "##Optimization",
        "Image lazy loading for 42 destination galleries",
        "Code splitting for admin vs user views",
        "CDN (Vercel Edge) for static assets",
        "PostgreSQL query optimization with proper indexing",
        ">>PWA-first approach = 60% lower development cost vs native apps",
    ], BRAND_PURPLE)

    make_content_slide(prs, "MCP Server Integration", [
        "##Model Context Protocol (MCP) Server",
        "AI agent integration for intelligent trip recommendations",
        "Built with @modelcontextprotocol/sdk",
        "##Capabilities",
        "Natural language trip search: 'Find me a beach trip under ₹10K'",
        "Group preference matching: Analyze interests, suggest optimal trip",
        "Dynamic pricing suggestions based on demand/seasonality",
        "##Architecture",
        "MCP server runs as standalone service alongside main API",
        "Exposes tools: search_outings, recommend_trip, check_availability",
        "Integrates with LLM providers for natural language understanding",
        ">>AI-powered recommendations increase conversion by estimated 25%",
    ], BRAND_PURPLE)

    make_content_slide(prs, "DevOps & Infrastructure", [
        "##Current Infrastructure",
        "Render.com: Auto-deploy from GitHub, managed PostgreSQL, SSL",
        "Vercel: Edge deployment for frontend, instant rollbacks",
        "Docker: Containerized backend for consistent environments",
        "##CI/CD Pipeline",
        "GitHub → Auto-deploy to Render (backend) + Vercel (frontend)",
        "Test suite: Automated tests in tests/test_all.js",
        "##Monitoring & Observability",
        "Application logs via Render dashboard",
        "Error tracking: Structured error responses with status codes",
        "Database monitoring: Connection pool health, query performance",
        "##Scale-Up Plan",
        "Horizontal scaling: Render auto-scale on traffic spikes",
        "Database: Read replicas for analytics queries",
        "CDN: Cloudflare for global edge caching",
        ">>Current architecture handles 10,000+ concurrent users",
    ], BRAND_PURPLE)

    make_content_slide(prs, "Technical Roadmap", [
        "##Phase 1: Current (MVP Live)",
        "Full-featured PWA, Razorpay payments, Admin dashboard, MCP server",
        "##Phase 2: Q3-Q4 2026",
        "React Native mobile app (iOS + Android)",
        "Real-time notifications (WebSocket/SSE)",
        "AI recommendation engine v2 (collaborative filtering)",
        "##Phase 3: 2027",
        "Organizer self-service portal with API",
        "Advanced analytics dashboard (Metabase/custom)",
        "Multi-region database deployment",
        "##Phase 4: 2028+",
        "Microservices architecture migration",
        "ML-powered dynamic pricing",
        "Social features: User profiles, connections, trip albums",
        ">>All phases maintain backward compatibility — zero downtime migrations",
    ], BRAND_PURPLE)

    make_closing_slide(prs, "Tech Lead")
    prs.save(os.path.join(OUTPUT_DIR, "07_Tech_Lead_Deck.pptx"))
    print("✅ Tech Lead deck created")


# ═══════════════════════════════════════════════════════════════
# 8. BUSINESS DEVELOPMENT
# ═══════════════════════════════════════════════════════════════
def create_business_development_ppt():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    make_title_slide(prs, "VIBES@Outing", "Business Development Deck", "Partnerships | Revenue Growth | Market Expansion | Strategy")

    make_content_slide(prs, "BD Vision & Strategy", [
        "##Mission",
        "Build a partner ecosystem that powers premium GenZ group outings at scale",
        "Transform VIBES@Outing from a marketplace into a platform",
        "##Strategic Pillars",
        "Resort & hospitality partnerships (supply side)",
        "Corporate & institutional partnerships (demand side)",
        "City expansion through local partner networks",
        "Revenue diversification through new business lines",
        ">>Every partnership should increase supply, demand, or both",
    ], BRAND_TEAL)

    make_content_slide(prs, "Partner Ecosystem", [
        "##Supply Partners",
        "Premium Resorts: 42 destinations, Tier-1 properties — bulk rate agreements",
        "Cab/Transport: Tempo Traveller & mini-bus fleet partners (3+ per city)",
        "Activity Vendors: Certified adventure operators (rafting, trekking, etc.)",
        "Local Guides: Heritage & wildlife experts at each destination",
        "##Demand Partners",
        "Colleges: Campus ambassadors at 50+ colleges in Bangalore",
        "Startups: GenZ team outing packages for 500+ Bangalore startups",
        "Co-working Spaces: WeWork, 91springboard, Cowrks — outing promotions",
        "##Strategic Partners",
        "Travel gear brands: Wildcraft, Decathlon — co-branded content + offers",
        "Food/Beverage: Cafe partnerships at trip destinations",
        "Insurance: Group travel insurance partnerships",
        ">>Goal: 200+ active partners across all categories by Year 1 end",
    ], BRAND_TEAL)

    make_table_slide(prs, "Revenue Opportunity by Partner Type", 
        ["Partner Type", "Deal Structure", "Revenue/Partner/Month", "Year 1 Target"],
        [
            ["Premium Resorts", "Negotiated bulk rates (30-40% off)", "₹50K-₹2L", "42 resorts"],
            ["Corporate Outings", "₹8K-₹15K per head, 15-50 pax", "₹2L-₹5L", "20 corporates"],
            ["College Tie-ups", "₹200/student commission", "₹30K-₹80K", "30 colleges"],
            ["Featured Listings", "₹5K-₹20K/month for visibility", "₹10K avg", "15 featured"],
            ["Affiliate/Brand Deals", "Revenue share model", "₹20K-₹50K", "10 brands"],
            ["Organizer Platform", "20% commission on bookings", "₹40K avg", "Phase 2"],
        ], BRAND_TEAL)

    make_content_slide(prs, "Corporate & Institutional Sales", [
        "##Corporate GenZ Packages",
        "Target: Startups & tech companies with young teams (avg age < 30)",
        "Customized team outing packages: Team building + adventure + resort",
        "Price: ₹8,000-₹15,000 per person (premium margin)",
        "##Value Proposition for Companies",
        "Zero coordination: We handle everything end-to-end",
        "Team bonding: Curated activities designed for team building",
        "Professional coordination: Dedicated trip manager for corporate groups",
        "##Sales Pipeline",
        "Target 500+ Bangalore startups via LinkedIn, startup events, co-working spaces",
        "Partnership with HR platforms: Keka, Darwinbox, Zoho People",
        "##Institutional Sales",
        "College fest sponsorships: Travel prize giveaways",
        "Alumni association partnerships: Reunion trip packages",
        ">>Corporate packages = highest margin business line (30-40% margin)",
    ], BRAND_TEAL)

    make_content_slide(prs, "City Expansion Strategy", [
        "##Expansion Roadmap",
        "City 1 (Current): Bangalore — 42 outings, full operations",
        "City 2 (Month 4): Hyderabad — 20+ outings, leverage AP/Telangana destinations",
        "City 3 (Month 5): Chennai — 20+ outings, Tamil Nadu coastline & hills",
        "City 4 (Month 7): Pune — 20+ outings, Western Ghats & Goa corridor",
        "City 5 (Month 9): Mumbai — 25+ outings, Konkan coast & Maharashtra forts",
        "##City Launch Playbook",
        "Step 1: Onboard 15+ resort partners in region (8 weeks)",
        "Step 2: Curate 20+ outings with local flavor (4 weeks)",
        "Step 3: Hire regional operations lead + 2 trip coordinators",
        "Step 4: Launch with college + startup partnerships",
        "Step 5: Scale through referral program + local influencers",
        ">>Each city reaches break-even within 3 months of launch",
    ], BRAND_TEAL)

    make_kpi_slide(prs, "Business Development KPIs", [
        ("Active Partners", "200+"),
        ("Corporate Clients (Y1)", "20"),
        ("College Tie-ups", "30+"),
        ("Cities Operational", "5"),
        ("Partner Satisfaction", "> 4.5/5"),
        ("BD Revenue (Y1)", "₹28L+"),
    ])

    make_content_slide(prs, "Strategic Partnerships & Alliances", [
        "##Travel Industry Alliances",
        "Karnataka Tourism Board: Official GenZ travel partner (explore)",
        "IRCTC: Curated train + outing combo packages for long-distance trips",
        "OYO/Treebo: Bulk inventory access for budget-premium destinations",
        "##Technology Partnerships",
        "Razorpay: Preferred payment partner — co-marketing opportunities",
        "Google Maps Platform: Enhanced trip routing & destination info",
        "WhatsApp Business API: Automated trip notifications & group management",
        "##Brand Collaborations",
        "Wildcraft/Decathlon: Co-branded travel kits for trip participants",
        "Red Bull/Monster: Adventure activity sponsorship",
        "Spotify: Curated trip playlists — 'VIBES@Outing Road Trip Mix'",
        ">>Strategic partnerships amplify reach without proportional cost increase",
    ], BRAND_TEAL)

    make_content_slide(prs, "Organizer Platform (Phase 2)", [
        "##Vision: Two-Sided Marketplace",
        "Allow verified organizers to create and manage their own trips",
        "VIBES@Outing provides: Platform, payments, trust layer, marketing",
        "##Organizer Benefits",
        "Ready-made audience of GenZ travelers",
        "Payment collection via Razorpay (no setup needed)",
        "ID verification for safety",
        "Ratings & reviews for credibility",
        "##Revenue Model",
        "20% commission on all organizer bookings",
        "Featured listing: ₹5K-₹20K/month for premium visibility",
        "##Quality Control",
        "Organizer verification: Business registration, reviews, background check",
        "Trip quality audit: First 3 trips supervised by VIBES@Outing coordinator",
        ">>Organizer platform transforms unit economics — 10x supply with minimal ops cost",
    ], BRAND_TEAL)

    make_closing_slide(prs, "Business Development")
    prs.save(os.path.join(OUTPUT_DIR, "08_Business_Development_Deck.pptx"))
    print("✅ Business Development deck created")


# ═══════════════════════════════════════════════════════════════
# RUN ALL
# ═══════════════════════════════════════════════════════════════
if __name__ == "__main__":
    print("🚀 Generating VIBES@Outing role-specific presentations...\n")
    
    create_founder_ceo_ppt()
    create_marketing_manager_ppt()
    create_digital_media_ppt()
    create_operations_ppt()
    create_finance_manager_ppt()
    create_customer_relationship_ppt()
    create_tech_lead_ppt()
    create_business_development_ppt()
    
    print(f"\n✨ All 8 presentations saved to: {OUTPUT_DIR}")
    print("Files created:")
    for f in sorted(os.listdir(OUTPUT_DIR)):
        if f.endswith('.pptx'):
            print(f"  📊 {f}")
